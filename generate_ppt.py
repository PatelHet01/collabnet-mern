from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
title_only_layout = prs.slide_layouts[5]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "CollabNet: Influencer Marketing Platform"
subtitle.text = "Connecting Forward-Thinking Brands with Top-Tier Creators\n\nMERN Stack Web Application"

# Slide 2: Problem Statement
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "The Problem"
tf = body_shape.text_frame
tf.text = "Traditional influencer marketing is fragmented."
p = tf.add_paragraph()
p.text = "Brands struggle to discover authentic creators."
p.level = 1
p = tf.add_paragraph()
p.text = "Creators face delayed payments and unclear project scopes."
p.level = 1
p = tf.add_paragraph()
p.text = "Agencies lack tools to manage large rosters."
p.level = 1

# Slide 3: The Solution (CollabNet)
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "The Solution: CollabNet"
tf = body_shape.text_frame
tf.text = "A centralized, Role-Based Access Control (RBAC) platform."
p = tf.add_paragraph()
p.text = "Streamlined Campaign Workflows."
p.level = 1
p = tf.add_paragraph()
p.text = "Transparent matching between Brands and Creators."
p.level = 1
p = tf.add_paragraph()
p.text = "Real-time tracking of deliverables and payments."
p.level = 1

# Slide 4: Key Features
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Key Features"
tf = body_shape.text_frame
tf.text = "Multi-Dashboard Architecture"
p = tf.add_paragraph()
p.text = "Brand Dashboard (Campaign Creation, Proposals)"
p.level = 1
p = tf.add_paragraph()
p.text = "Creator Studio (Earnings, Job Browsing, Chats)"
p.level = 1
p = tf.add_paragraph()
p.text = "Agency Hub (Roster Management)"
p.level = 1
p = tf.add_paragraph()
p.text = "Super Admin Command Center (Platform health)"
p.level = 1

# Slide 5: Tech Stack
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Technology Stack"
tf = body_shape.text_frame
tf.text = "Built on the MERN Stack for high performance."
p = tf.add_paragraph()
p.text = "Frontend: React.js, Vite, React Router, Vanilla CSS"
p.level = 1
p = tf.add_paragraph()
p.text = "Backend: Node.js, Express.js"
p.level = 1
p = tf.add_paragraph()
p.text = "Database: MongoDB (Configured for future scaling)"
p.level = 1
p = tf.add_paragraph()
p.text = "UI Design: Minimalist Black & White Theme"
p.level = 1

# Slide 6: Future Roadmap
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Future Roadmap"
tf = body_shape.text_frame
tf.text = "Phase 2 & 3 Enhancements"
p = tf.add_paragraph()
p.text = "Stripe API Integration for automated payouts."
p.level = 1
p = tf.add_paragraph()
p.text = "Social Media API Integrations (TikTok, IG) for live engagement data."
p.level = 1
p = tf.add_paragraph()
p.text = "AI-Driven Influencer Matching."
p.level = 1

# Slide 7: Conclusion
slide = prs.slides.add_slide(title_only_layout)
title_shape = slide.shapes.title
title_shape.text = "Thank You\n\nAny Questions?"

# Save presentation
prs.save("CollabNet_Presentation.pptx")
print("Presentation generated successfully!")
